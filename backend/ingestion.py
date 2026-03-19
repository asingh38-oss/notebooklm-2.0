"""
backend/ingestion.py
--------------------
Ingestion pipeline: extract → chunk → embed → store in ChromaDB.

Chunking strategies:
  - "fixed"     : fixed character windows with overlap
  - "sentence"  : sentence-boundary-aware sliding window
  - "recursive" : recursive character splitter (default, best quality)
"""

import hashlib
import os
import re
from pathlib import Path
from typing import Literal, Optional
from urllib.parse import urlparse

import chromadb
from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction

from backend import storage

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

OPENAI_API_KEY: str  = os.getenv("OPENAI_API_KEY", "")
CHUNK_SIZE: int      = int(os.getenv("CHUNK_SIZE", "800"))
CHUNK_OVERLAP: int   = int(os.getenv("CHUNK_OVERLAP", "150"))
EMBED_MODEL: str     = os.getenv("EMBED_MODEL", "text-embedding-3-small")

ChunkStrategy = Literal["fixed", "sentence", "recursive"]
DEFAULT_STRATEGY: ChunkStrategy = "recursive"


# ---------------------------------------------------------------------------
# ChromaDB helpers
# ---------------------------------------------------------------------------

def _embedding_fn():
    return OpenAIEmbeddingFunction(
        api_key=OPENAI_API_KEY,
        model_name=EMBED_MODEL,
    )


def _get_collection(username: str, notebook_id: str):
    chroma_path = storage.get_chroma_path(username, notebook_id)
    client = chromadb.PersistentClient(path=chroma_path)
    return client.get_or_create_collection(
        name="sources",
        embedding_function=_embedding_fn(),
        metadata={"hnsw:space": "cosine"},
    )


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def _extract_pdf(file_path: Path) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(str(file_path)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    return "\n\n".join(texts)


def _extract_pptx(file_path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        lines.append(line)
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _extract_url(url: str) -> str:
    import trafilatura
    downloaded = trafilatura.fetch_url(url)
    if downloaded is None:
        raise ValueError(f"Could not fetch URL: {url}")
    text = trafilatura.extract(downloaded, include_links=False, include_tables=True)
    if not text:
        raise ValueError(f"No readable text found at: {url}")
    return text


def extract_text(source: str | Path) -> tuple[str, str]:
    src = str(source)
    parsed = urlparse(src)

    if parsed.scheme in ("http", "https"):
        text = _extract_url(src)
        name = parsed.netloc + parsed.path.replace("/", "_")
        return text, name

    p = Path(src)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(p), p.name
    elif suffix in (".pptx", ".ppt"):
        return _extract_pptx(p), p.name
    elif suffix in (".txt", ".md"):
        return p.read_text(encoding="utf-8", errors="ignore"), p.name
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


# ---------------------------------------------------------------------------
# Chunking strategies
# ---------------------------------------------------------------------------

def _fixed_chunks(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start: start + size])
        start += size - overlap
    return [c.strip() for c in chunks if c.strip()]


def _sentence_chunks(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks, current, current_len = [], [], 0
    for sent in sentences:
        slen = len(sent)
        if current_len + slen > size and current:
            chunks.append(" ".join(current))
            tail, tlen = [], 0
            for s in reversed(current):
                tlen += len(s)
                tail.insert(0, s)
                if tlen >= overlap:
                    break
            current, current_len = tail, sum(len(s) for s in tail)
        current.append(sent)
        current_len += slen
    if current:
        chunks.append(" ".join(current))
    return [c.strip() for c in chunks if c.strip()]


def _recursive_chunks(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    separators = ["\n\n", "\n", ". ", " ", ""]

    def _split(t: str, seps: list[str]) -> list[str]:
        if len(t) <= size:
            return [t] if t.strip() else []
        if not seps:
            return _fixed_chunks(t, size, overlap)
        sep = seps[0]
        parts = t.split(sep)
        result, current = [], ""
        for part in parts:
            candidate = current + (sep if current else "") + part
            if len(candidate) <= size:
                current = candidate
            else:
                if current:
                    result.extend(_split(current, seps[1:]))
                current = part
        if current:
            result.extend(_split(current, seps[1:]))
        return result

    raw = _split(text, separators)
    overlapped = []
    for i, chunk in enumerate(raw):
        if i > 0 and overlap > 0:
            chunk = raw[i - 1][-overlap:] + " " + chunk
        overlapped.append(chunk.strip())
    return [c for c in overlapped if c]


def chunk_text(text: str, strategy: ChunkStrategy = DEFAULT_STRATEGY) -> list[str]:
    if strategy == "fixed":
        return _fixed_chunks(text)
    elif strategy == "sentence":
        return _sentence_chunks(text)
    else:
        return _recursive_chunks(text)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def ingest_source(
    username: str,
    notebook_id: str,
    source: str | Path,
    chunk_strategy: ChunkStrategy = DEFAULT_STRATEGY,
    raw_bytes: Optional[bytes] = None,
) -> dict:
    text, source_name = extract_text(source)

    if raw_bytes is not None:
        storage.save_raw_file(username, notebook_id, source_name, raw_bytes)

    storage.save_extracted_text(username, notebook_id, source_name, text)

    chunks = chunk_text(text, strategy=chunk_strategy)

    ids, documents, metadatas = [], [], []
    for i, chunk in enumerate(chunks):
        chunk_id = hashlib.md5(f"{source_name}_{i}_{chunk[:40]}".encode()).hexdigest()
        ids.append(chunk_id)
        documents.append(chunk)
        metadatas.append({
            "source": source_name,
            "chunk_index": i,
            "total_chunks": len(chunks),
            "strategy": chunk_strategy,
        })

    collection = _get_collection(username, notebook_id)
    for start in range(0, len(ids), 100):
        collection.upsert(
            ids=ids[start: start + 100],
            documents=documents[start: start + 100],
            metadatas=metadatas[start: start + 100],
        )

    return {
        "source_name": source_name,
        "char_count": len(text),
        "chunk_count": len(chunks),
        "strategy": chunk_strategy,
    }


def delete_source(username: str, notebook_id: str, source_name: str) -> int:
    collection = _get_collection(username, notebook_id)
    results = collection.get(where={"source": source_name})
    if results and results["ids"]:
        collection.delete(ids=results["ids"])
        return len(results["ids"])
    return 0


def list_indexed_sources(username: str, notebook_id: str) -> list[str]:
    collection = _get_collection(username, notebook_id)
    results = collection.get(include=["metadatas"])
    if not results["metadatas"]:
        return []
    return sorted({m["source"] for m in results["metadatas"]})